"""Paper / PPT text extraction from PDF, DOCX, PPTX, or TXT.

Limitations:
- Only extracts text; images, charts, formulas, SmartArt, and diagrams are NOT understood.
- Scanned/image-based PDFs will produce empty or garbled output.
- Two-column layouts may have reading order issues.
- Tables may lose structure.

For full paper understanding, students should ensure key content exists as
searchable text in the uploaded file. Image-heavy PPTs should be supplemented
with a text-based PDF version.
"""
import io
from pypdf import PdfReader
from docx import Document


def extract_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            parts.append(f"[第 {i+1} 页]\n{text}")
    return "\n\n".join(parts).strip()


def extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs).strip()


def extract_pptx(data: bytes) -> str:
    """Extract text from PPTX slides. Does NOT understand images, charts, or diagrams."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text.strip())
        if slide_texts:
            texts.append(f"[幻灯片 {i+1}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def extract(filename: str, data: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return extract_pdf(data)
    if name.endswith(".docx"):
        return extract_docx(data)
    if name.endswith(".pptx"):
        return extract_pptx(data)
    if name.endswith(".txt"):
        return data.decode("utf-8", errors="ignore")
    raise ValueError(f"不支持的文件格式: {filename}。请上传 PDF、DOCX、PPTX 或 TXT 文件。")
